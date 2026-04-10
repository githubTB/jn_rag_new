from __future__ import annotations

from collections.abc import Iterable
from dataclasses import dataclass
import os
import tempfile
from pathlib import Path

from .base import BaseExtractor
from models.document import Document

_MIN_PICTURE_OCR_DIM = 900000
_FLOW_NODE_MIN_TEXT = 8
_FLOW_Y_TOLERANCE = 500000


@dataclass
class _SlideEntry:
    top: int
    left: int
    order: int
    text: str


class PptxExtractor(BaseExtractor):
    """Extract text from .pptx files, one Document per slide."""

    def __init__(self, file_path: str):
        self._file_path = file_path

    def extract(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self._file_path)
        documents: list[Document] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            entries = self._extract_slide_entries(slide.shapes)
            lines = [entry.text for entry in entries]
            flow_lines = self._build_flow_lines(entries)
            if not lines and not flow_lines:
                continue
            if lines:
                documents.append(
                    Document(
                        page_content="\n".join(lines),
                        metadata={
                            "source": self._file_path,
                            "slide": slide_num,
                            "label": "slide",
                            "flow_lines": flow_lines,
                        },
                    )
                )
            for flow_index, flow_line in enumerate(flow_lines):
                documents.append(
                    Document(
                        page_content=flow_line,
                        metadata={
                            "source": self._file_path,
                            "slide": slide_num,
                            "label": "flow",
                            "flow_index": flow_index,
                            "flow_lines": flow_lines,
                        },
                    )
                )

        return documents

    def _extract_slide_entries(self, shapes: Iterable) -> list[_SlideEntry]:
        entries: list[_SlideEntry] = []
        seen: set[str] = set()
        self._collect_shape_entries(shapes, entries, seen)
        entries.sort(key=lambda item: (item.top, item.left, item.order))
        return entries

    def _collect_shape_entries(
        self,
        shapes: Iterable,
        entries: list[_SlideEntry],
        seen: set[str],
        *,
        parent_top: int = 0,
        parent_left: int = 0,
        depth: int = 0,
    ) -> None:
        ordered_shapes = sorted(
            list(shapes),
            key=lambda shape: (
                int(getattr(shape, "top", 0) or 0),
                int(getattr(shape, "left", 0) or 0),
            ),
        )

        for index, shape in enumerate(ordered_shapes):
            shape_top = parent_top + int(getattr(shape, "top", 0) or 0)
            shape_left = parent_left + int(getattr(shape, "left", 0) or 0)

            if hasattr(shape, "shapes"):
                self._collect_shape_entries(
                    shape.shapes,
                    entries,
                    seen,
                    parent_top=shape_top,
                    parent_left=shape_left,
                    depth=depth + 1,
                )

            text_lines = self._extract_text_lines(shape)
            for line_idx, line in enumerate(text_lines):
                normalized = line.strip()
                if not normalized or normalized in seen:
                    continue
                seen.add(normalized)
                entries.append(_SlideEntry(shape_top, shape_left, index * 100 + line_idx + depth * 10000, normalized))

            if getattr(shape, "has_table", False):
                for row_idx, row_text in enumerate(self._extract_table_lines(shape)):
                    if not row_text or row_text in seen:
                        continue
                    seen.add(row_text)
                    entries.append(_SlideEntry(shape_top, shape_left, index * 100 + row_idx + 50000, row_text))

            for ocr_idx, ocr_text in enumerate(self._extract_picture_ocr_lines(shape)):
                if not ocr_text or ocr_text in seen:
                    continue
                seen.add(ocr_text)
                entries.append(_SlideEntry(shape_top, shape_left, index * 100 + ocr_idx + 80000, ocr_text))

    def _extract_text_lines(self, shape) -> list[str]:
        if not getattr(shape, "has_text_frame", False):
            return []
        lines: list[str] = []
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                lines.append(line)
        return lines

    def _extract_table_lines(self, shape) -> list[str]:
        table = shape.table
        lines: list[str] = []
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                lines.append(" | ".join(row_texts))
        return lines

    def _extract_picture_ocr_lines(self, shape) -> list[str]:
        if not hasattr(shape, "image"):
            return []

        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        if width < _MIN_PICTURE_OCR_DIM or height < _MIN_PICTURE_OCR_DIM:
            return []

        image = shape.image
        ext = image.ext or "png"
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(image.blob)
                tmp_path = tmp.name

            from extractor.ocr_router import route_ocr

            docs = route_ocr(tmp_path, doc_type="unknown")
            lines: list[str] = []
            for doc in docs:
                for line in (doc.page_content or "").splitlines():
                    normalized = line.strip()
                    if normalized:
                        lines.append(normalized)
            return lines
        except Exception:
            return []
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

    def _build_flow_lines(self, entries: list[_SlideEntry]) -> list[str]:
        node_entries = [
            entry for entry in entries
            if entry.text
            and len(entry.text) <= _FLOW_NODE_MIN_TEXT
            and "整体" not in entry.text
            and "步骤" not in entry.text
        ]
        if len(node_entries) < 3:
            return []

        rows: list[list[_SlideEntry]] = []
        for entry in sorted(node_entries, key=lambda item: (item.top, item.left)):
            matched = False
            for row in rows:
                if abs(row[0].top - entry.top) <= _FLOW_Y_TOLERANCE:
                    row.append(entry)
                    matched = True
                    break
            if not matched:
                rows.append([entry])

        rows = [sorted(row, key=lambda item: item.left) for row in rows if len(row) >= 2]
        if not rows:
            return []

        rows.sort(key=lambda row: (-len(row), row[0].top))
        main_row = rows[0]
        flow_lines = [f"流程主线: {' -> '.join(entry.text for entry in main_row)}"]

        branch_lines: list[str] = []
        main_lefts = {entry.left for entry in main_row}
        for row in rows[1:]:
            shared = [entry for entry in row if any(abs(entry.left - left) <= 200000 for left in main_lefts)]
            if shared:
                branch_lines.append(f"流程分支: {' -> '.join(entry.text for entry in row)}")

        return flow_lines + branch_lines
